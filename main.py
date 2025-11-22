import asyncio
import sys
import os
import tempfile
from typing import List
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import speech_recognition as sr
from pydub import AudioSegment
from dotenv import load_dotenv
from fastapi.staticfiles import StaticFiles
import docx
import pptx

# Windows compatibility fix
if sys.platform == "win32":
    asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

# Import for ElevenLabs - NEW VERSION
from elevenlabs.client import ElevenLabs

# ---------------- Load ENV & Setup ---------------- #
load_dotenv()

# Get API keys from environment
ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

if not ELEVENLABS_API_KEY:
    raise ValueError("Please set the ELEVENLABS_API_KEY environment variable.")
if not GROQ_API_KEY:
    raise ValueError("Please set the GROQ_API_KEY environment variable.")

# Configure the ElevenLabs client - NEW METHOD
elevenlabs_client = ElevenLabs(api_key=ELEVENLABS_API_KEY)

# Configure Groq client
from openai import OpenAI
groq_client = OpenAI(api_key=GROQ_API_KEY, base_url="https://api.groq.com/openai/v1")

# ---------------- FastAPI App Setup ---------------- #
app = FastAPI(title="RAG + Voice QA with ElevenLabs TTS")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------- Core RAG Variables ---------------- #
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
dimension = 384
index = faiss.IndexFlatL2(dimension)
documents: List[str] = []

# ---------------- Response Models ---------------- #
class UploadResponse(BaseModel):
    message: str

class HealthResponse(BaseModel):
    status: str
    message: str

# ---------------- Document Processing Functions ---------------- #
def read_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[-1].lower()
    text = ""
    try:
        if ext == ".pdf":
            import fitz
            with fitz.open(file_path) as doc:
                text = " ".join([page.get_text() for page in doc])
        elif ext == ".docx":
            import docx
            d = docx.Document(file_path)
            text = " ".join([p.text for p in d.paragraphs])
        elif ext == ".pptx":
            import pptx
            prs = pptx.Presentation(file_path)
            text = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error reading file: {str(e)}")
    return text

def embed_and_store(text: str):
    global documents, index
    # Split text into chunks of 500 characters with overlap
    chunk_size = 500
    overlap = 50
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap  # Overlap chunks for better context
    
    if chunks:
        embeddings = embedding_model.encode(chunks)
        index.add(np.array(embeddings, dtype=np.float32))
        documents.extend(chunks)

def retrieve_context(query: str, top_k: int = 3) -> str:
    if not documents:
        return "No documents available. Please upload documents first."
    
    query_emb = embedding_model.encode([query])
    D, I = index.search(np.array(query_emb, dtype=np.float32), top_k)
    
    # Retrieve the most relevant chunks
    relevant_chunks = []
    for i in I[0]:
        if i < len(documents):
            relevant_chunks.append(documents[i])
    
    return " ".join(relevant_chunks)

# ---------------- Routes ---------------- #
@app.get("/")
async def root():
    return {"message": "RAG Voice QA API is running!"}

@app.get("/health", response_model=HealthResponse)
async def health_check():
    return {
        "status": "healthy",
        "message": "API is running successfully with Llama 3.1 and ElevenLabs integration"
    }

@app.get("/voices")
async def list_voices():
    """Get list of available ElevenLabs voices"""
    try:
        voices = elevenlabs_client.voices.get_all()
        voice_list = []
        for voice in voices.voices:
            voice_list.append({
                "id": voice.voice_id,
                "name": voice.name,
                "category": voice.category
            })
        return voice_list
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching voices: {str(e)}")

@app.post("/upload_doc", response_model=UploadResponse)
async def upload_doc(file: UploadFile = File(...)):
    try:
        # Validate file type
        allowed_extensions = {'.pdf', '.docx', '.pptx', '.txt'}
        file_extension = os.path.splitext(file.filename)[-1].lower()
        if file_extension not in allowed_extensions:
            raise HTTPException(status_code=400, detail="Unsupported file format. Please upload PDF, DOCX, PPTX, or TXT files.")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp:
            content = await file.read()
            if not content:
                raise HTTPException(status_code=400, detail="Uploaded file is empty.")
            tmp.write(content)
            tmp_path = tmp.name

        text = read_file(tmp_path)
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the document.")
        
        embed_and_store(text)
        os.remove(tmp_path)
        
        return {"message": f"Document '{file.filename}' uploaded and processed successfully. {len(documents)} chunks added."}
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")

@app.post("/ask_audio")
async def ask_audio(file: UploadFile = File(...)):
    try:
        # Validate audio file
        if not file.filename.lower().endswith(('.wav', '.mp3', '.m4a', '.flac')):
            raise HTTPException(status_code=400, detail="Unsupported audio format. Please upload WAV, MP3, M4A, or FLAC files.")
        
        # --- STEP 1: Speech-to-Text --- #
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[-1]) as tmp:
            content = await file.read()
            if not content:
                raise HTTPException(status_code=400, detail="Uploaded audio file is empty.")
            tmp.write(content)
            tmp_path = tmp.name

        try:
            # Convert to WAV if needed
            sound = AudioSegment.from_file(tmp_path)
            wav_path = tmp_path + "_converted.wav"
            sound.export(wav_path, format="wav")

            # Speech recognition
            recognizer = sr.Recognizer()
            with sr.AudioFile(wav_path) as source:
                # Adjust for ambient noise
                recognizer.adjust_for_ambient_noise(source, duration=0.5)
                audio_data = recognizer.record(source)
                query = recognizer.recognize_google(audio_data)
                
        except sr.UnknownValueError:
            raise HTTPException(status_code=400, detail="Could not understand the audio. Please speak clearly.")
        except sr.RequestError as e:
            raise HTTPException(status_code=500, detail=f"Speech recognition service error: {str(e)}")
        finally:
            # Clean up temporary files
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
            if os.path.exists(wav_path):
                os.remove(wav_path)

        if not query.strip():
            raise HTTPException(status_code=400, detail="No speech detected in the audio file.")

        # --- STEP 2: Retrieve Context --- #
        context = retrieve_context(query)
        
        # --- STEP 3: Generate Answer using Llama 3.1 --- #
        prompt = f"""Based on the following context, please provide a clear and concise answer to the question.

CONTEXT:
{context}

QUESTION:
{query}

Please provide a direct answer in 25-35 words based solely on the context above. If the context doesn't contain the answer, say "I cannot find the answer in the provided documents."
ANSWER:"""
        
        try:
            response = groq_client.chat.completions.create(
                model="llama-3.1-8b-instant",  # Using Llama 3.1
                messages=[{"role": "user", "content": prompt}],
                max_tokens=100,
                temperature=0.3  # Lower temperature for more factual responses
            )
            answer = response.choices[0].message.content.strip()
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error generating answer: {str(e)}")

        # --- STEP 4: Convert the text answer to speech with ElevenLabs --- #
        try:
            # Get available voices and use Rachel if available
            voices = elevenlabs_client.voices.get_all()
            voice_id = None
            
            # Try to find Rachel voice
            for voice in voices.voices:
                if voice.name.lower() == "rachel":
                    voice_id = voice.voice_id
                    break
            
            # If Rachel not found, use the first available voice
            if not voice_id and voices.voices:
                voice_id = voices.voices[0].voice_id
            
            if not voice_id:
                raise HTTPException(status_code=500, detail="No voices available in ElevenLabs account")
            
            # Generate speech - NEW METHOD for ElevenLabs v1.x
            audio_response = elevenlabs_client.text_to_speech.convert(
                voice_id=voice_id,
                model_id="eleven_monolingual_v1",
                text=answer,
                output_format="mp3_44100_128"
            )
            
            # Collect all audio chunks
            audio_content = b"".join([chunk for chunk in audio_response])
            
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error generating speech: {str(e)}")

        # --- STEP 5: Stream the audio back to the client --- #
        return StreamingResponse(
            iter([audio_content]),
            media_type="audio/mpeg",
            headers={
                "Content-Disposition": "attachment; filename=answer.mp3",
                "Transcript-Text": answer,
                "X-Original-Query": query
            }
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")

@app.get("/documents/status")
async def documents_status():
    return {
        "total_documents": len(documents),
        "total_chunks": len(documents),
        "index_size": index.ntotal if hasattr(index, 'ntotal') else 0
    }

@app.post("/clear_documents")
async def clear_documents():
    global documents, index
    documents.clear()
    index = faiss.IndexFlatL2(dimension)
    return {"message": "All documents and index have been cleared."}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)